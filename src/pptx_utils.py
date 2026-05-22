from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from .rules import EMU_PER_INCH


def emu_to_inches(value: int | float | None) -> float:
    if value is None:
        return 0.0
    return round(float(value) / EMU_PER_INCH, 3)


def shape_bbox(shape) -> dict[str, float]:
    return {
        "left": emu_to_inches(getattr(shape, "left", 0)),
        "top": emu_to_inches(getattr(shape, "top", 0)),
        "width": emu_to_inches(getattr(shape, "width", 0)),
        "height": emu_to_inches(getattr(shape, "height", 0)),
    }


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                yield subshape


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts = []
    for paragraph in shape.text_frame.paragraphs:
        para = "".join(run.text for run in paragraph.runs)
        if para.strip():
            parts.append(para)
    return "\n".join(parts).strip()


def slide_text(slide) -> str:
    texts = []
    for shape in iter_shapes(slide):
        txt = shape_text(shape)
        if txt:
            texts.append(txt)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return "\n".join(texts)


def get_title_shape(slide):
    # Prefer placeholder title.
    try:
        if slide.shapes.title and shape_text(slide.shapes.title).strip():
            return slide.shapes.title
    except Exception:
        pass

    # Fallback: top-most text shape with meaningful text.
    candidates = []
    for shape in iter_shapes(slide):
        txt = normalize_text(shape_text(shape))
        if txt and len(txt) > 4:
            candidates.append((getattr(shape, "top", 10**18), getattr(shape, "left", 10**18), shape))
    if not candidates:
        return None
    candidates.sort(key=lambda x: (x[0], x[1]))
    return candidates[0][2]


def get_text_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return []
    runs = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            runs.append((paragraph, run))
    return runs


def first_run_font(shape):
    runs = get_text_runs(shape)
    if not runs:
        return None
    return runs[0][1].font


def font_size_pt(font) -> float | None:
    try:
        if font and font.size:
            return round(font.size.pt, 1)
    except Exception:
        return None
    return None


def font_rgb_hex(font) -> str | None:
    try:
        color = font.color
        if color and color.rgb:
            return f"#{str(color.rgb).upper()}"
    except Exception:
        return None
    return None


def fill_rgb_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill and fill.fore_color and fill.fore_color.rgb:
            return f"#{str(fill.fore_color.rgb).upper()}"
    except Exception:
        return None
    return None


def table_text(shape) -> list[list[str]]:
    if not getattr(shape, "has_table", False):
        return []
    out = []
    for row in shape.table.rows:
        out.append([cell.text.strip() for cell in row.cells])
    return out


def extract_notes(slide) -> list[str]:
    try:
        notes = slide.notes_slide.notes_text_frame.text
    except Exception:
        return []
    lines = [normalize_text(x) for x in notes.splitlines()]
    return [x for x in lines if x and not re.fullmatch(r"\d+", x)]


def load_presentation(path: str | Path) -> Presentation:
    return Presentation(str(path))
