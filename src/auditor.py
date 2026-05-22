from __future__ import annotations

import re
import tempfile
from pathlib import Path
from typing import Iterable

from pptx import Presentation

from .models import AuditReport, Issue, SlideSummary
from .pptx_utils import (
    extract_notes,
    fill_rgb_hex,
    first_run_font,
    font_rgb_hex,
    font_size_pt,
    get_title_shape,
    iter_shapes,
    load_presentation,
    normalize_text,
    shape_bbox,
    shape_text,
    slide_text,
    table_text,
)
from .rules import (
    Enterprise_COLORS,
    Enterprise_CONFIDENTIALITY_TEXT,
    ENDING_PUNCTUATION,
    MIN_BODY_LEFT_INCHES,
    TERMINOLOGY_REPLACEMENTS,
)


def make_issue(slide_num: int, rule_id: str, category: str, severity: str, title: str,
               evidence: str, recommendation: str, confidence: float = 0.85, shape=None) -> Issue:
    return Issue(
        slide_number=slide_num,
        rule_id=rule_id,
        category=category,
        severity=severity,  # type: ignore[arg-type]
        title=title,
        evidence=evidence,
        recommendation=recommendation,
        confidence=confidence,
        shape_name=getattr(shape, "name", None) if shape is not None else None,
        bbox=shape_bbox(shape) if shape is not None else None,
    )


def check_title(slide, slide_num: int) -> list[Issue]:
    issues: list[Issue] = []
    title_shape = get_title_shape(slide)
    if title_shape is None:
        return [make_issue(slide_num, "TITLE_MISSING", "Slide title", "high", "Missing slide title", "No title-like text shape was detected.", "Add a clear, takeaway-driven slide title.", 0.8)]

    title = normalize_text(shape_text(title_shape))
    font = first_run_font(title_shape)
    size = font_size_pt(font) if font else None
    is_bold = bool(getattr(font, "bold", False)) if font else False

    if title.endswith(ENDING_PUNCTUATION):
        issues.append(make_issue(slide_num, "TITLE_END_PUNCT", "Slide title", "medium", "Title ends with punctuation", f"Title: {title}", "Remove ending punctuation from the headline.", 0.95, title_shape))

    # Approximate 3-line warning: explicit line breaks or long title.
    line_count = title.count("\n") + 1
    if line_count > 3 or len(title) > 145:
        issues.append(make_issue(slide_num, "TITLE_TOO_LONG", "Slide title", "medium", "Headline may exceed 3 lines", f"Title length: {len(title)} characters.", "Shorten the headline so it communicates the takeaway in no more than 3 lines.", 0.7, title_shape))

    if size is not None and abs(size - 24) > 1.0:
        issues.append(make_issue(slide_num, "TITLE_SIZE", "Slide title", "medium", "Title font size differs from 24 pt", f"Detected title size: {size} pt.", "Use Calibri Bold 24 pt for content slide headlines.", 0.7, title_shape))

    if font is not None and not is_bold and slide_num != 1:
        issues.append(make_issue(slide_num, "TITLE_BOLD", "Slide title", "low", "Title may not be bold", "The first title run was not detected as bold.", "Use Calibri Bold for content slide headlines.", 0.65, title_shape))

    # Topic-only heuristic.
    if len(title.split()) <= 5 and not any(v in title.lower() for v in ["recommend", "suggest", "increase", "decrease", "reflect", "support", "continue", "remove", "add", "stable", "evolved"]):
        issues.append(make_issue(slide_num, "TITLE_TAKEAWAY", "Writing style", "low", "Title may be topic-only", f"Title: {title}", "Use an insight-driven headline that communicates the takeaway, not just the topic.", 0.55, title_shape))

    return issues


def check_confidentiality(slide, slide_num: int) -> list[Issue]:
    txt = normalize_text(slide_text(slide))
    if Enterprise_CONFIDENTIALITY_TEXT.lower() not in txt.lower():
        return [make_issue(slide_num, "CONFIDENTIALITY_MISSING", "Confidentiality", "high", "Missing confidentiality statement", "Expected Enterprise confidentiality statement was not found on this slide.", "Add the standard confidentiality statement without paraphrasing.", 0.95)]
    return []


def check_title_slide(slide, slide_num: int) -> list[Issue]:
    if slide_num != 1:
        return []
    txt = slide_text(slide)
    issues = []
    if "DRAFT" in txt.upper():
        issues.append(make_issue(slide_num, "TITLE_DRAFT", "Title slide", "high", "Title slide contains DRAFT", "The title slide includes a DRAFT designator.", "Remove DRAFT or other designator from the title slide unless explicitly required.", 0.95))
    if not re.search(r"\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec|January|February|March|April|May|June|July|August|September|October|November|December)\b|\d{1,2}/\d{1,2}/\d{2,4}|\d{4}", txt):
        issues.append(make_issue(slide_num, "TITLE_DATE_MISSING", "Title slide", "medium", "Title slide may be missing a valid date", "No obvious date was detected on the title slide.", "Add a valid date to the title slide.", 0.65))
    return issues


def check_layout_bounds(slide, slide_num: int) -> list[Issue]:
    issues = []
    title_shape = get_title_shape(slide)
    if title_shape is None:
        return issues
    title_box = shape_bbox(title_shape)
    title_left = title_box["left"]
    title_right = title_box["left"] + title_box["width"]

    for shape in iter_shapes(slide):
        if shape is title_shape:
            continue
        txt = normalize_text(shape_text(shape))
        if not txt and not getattr(shape, "has_table", False):
            continue
        box = shape_bbox(shape)
        if box["top"] < title_box["top"] + 0.2:
            continue
        if box["left"] < MIN_BODY_LEFT_INCHES:
            issues.append(make_issue(slide_num, "CONTENT_LEFT_MARGIN", "Layout", "medium", "Content is too close to left edge", f"Detected left position: {box['left']} in.", "Keep body content at least 1 inch from the left edge and aligned with the title/body grid.", 0.75, shape))
        if box["left"] < title_left - 0.15:
            issues.append(make_issue(slide_num, "CONTENT_LEFT_OF_TITLE", "Layout", "low", "Content starts left of the slide title", f"Content left: {box['left']} in; title left: {title_left} in.", "Align content with the title/body grid.", 0.65, shape))
        if box["left"] + box["width"] > title_right + 0.3:
            issues.append(make_issue(slide_num, "CONTENT_RIGHT_OF_TITLE", "Layout", "low", "Content may extend beyond title width", f"Content right: {round(box['left'] + box['width'], 2)} in; title right: {round(title_right, 2)} in.", "Keep content within the established title/body width.", 0.6, shape))
    return issues[:8]  # noise control


def check_bullets(slide, slide_num: int) -> list[Issue]:
    issues = []
    for shape in iter_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        for p in shape.text_frame.paragraphs:
            raw = "".join(run.text for run in p.runs).strip()
            if not raw or len(raw) < 4:
                continue
            is_bulletish = bool(getattr(p, "level", 0) > 0 or raw.startswith(("•", "▪", "-", "–", "➢")))
            if is_bulletish and raw.endswith(ENDING_PUNCTUATION):
                issues.append(make_issue(slide_num, "BULLET_END_PUNCT", "Bullets", "low", "Bullet ends with punctuation", raw[:180], "Remove ending punctuation from bullet text unless it is a full-sentence exception.", 0.65, shape))
    return issues[:6]


def check_terminology(slide, slide_num: int) -> list[Issue]:
    text = slide_text(slide)
    issues = []
    for bad, good in TERMINOLOGY_REPLACEMENTS.items():
        if re.search(rf"\b{re.escape(bad)}\b", text):
            issues.append(make_issue(slide_num, "TERMINOLOGY", "Terminology", "medium", f"Use '{good}' instead of '{bad}'", f"Found term: {bad}", f"Replace '{bad}' with '{good}' per Enterprise terminology guidelines.", 0.9))
    if "Company Name" in text:
        issues.append(make_issue(slide_num, "TABLE_COMPANY_NAME", "Tables", "low", "Table may use 'Company Name' instead of 'Company'", "Found table header: Company Name", "Consider using 'Company' if consistent with Enterprise table examples and annotated feedback.", 0.65))
    return issues


def check_tables(slide, slide_num: int) -> list[Issue]:
    issues = []
    text = slide_text(slide)
    has_money = bool(re.search(r"\$\d|\$\s*\d", text))
    has_table = any(getattr(s, "has_table", False) for s in iter_shapes(slide))
    if has_table and has_money and not re.search(r"\(\$\s*(?:MM|MMs|000s)\)|\$MM|\$000", text, flags=re.I):
        issues.append(make_issue(slide_num, "TABLE_SCALE_MISSING", "Tables", "medium", "Possible missing financial scale label", "Detected currency values but no obvious ($MM), ($MMs), or ($000s) scale label.", "Add a clear financial scale label to the table header or upper-left corner.", 0.55))

    for shape in iter_shapes(slide):
        if not getattr(shape, "has_table", False):
            continue
        rows = table_text(shape)
        flat = " ".join(cell for row in rows for cell in row)
        if "Incumbent" in flat:
            issues.append(make_issue(slide_num, "TABLE_EXECUTIVE_HEADER", "Tables", "medium", "Employee table may use 'Incumbent'", "Found 'Incumbent' in a table.", "Use 'Executive' for the column containing employee names.", 0.85, shape))
        # Heuristic: count/client row colors are hard to verify reliably in pptx, so flag only obvious text cases.
        if re.search(r"\bCount\b", flat) and "Client" not in flat:
            issues.append(make_issue(slide_num, "COUNT_ROW_REVIEW", "Tables", "low", "Review count/statistics row shading", "Found a Count row; statistics rows should use yellow shading, not client green.", "Verify statistics/count row uses Enterprise statistics-row shading (#FFFFDB).", 0.5, shape))
    return issues[:5]


def check_colors(slide, slide_num: int) -> list[Issue]:
    issues = []
    for shape in iter_shapes(slide):
        for color in [fill_rgb_hex(shape)]:
            if color and color.upper() not in Enterprise_COLORS:
                issues.append(make_issue(slide_num, "COLOR_OUTSIDE_PALETTE", "Colors", "low", "Shape color may be outside Enterprise palette", f"Detected color: {color}", "Use colors from the Enterprise color palette.", 0.6, shape))
        if getattr(shape, "has_text_frame", False):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    color = font_rgb_hex(run.font)
                    if color and color.upper() not in Enterprise_COLORS:
                        issues.append(make_issue(slide_num, "FONT_COLOR_OUTSIDE_PALETTE", "Colors", "low", "Font color may be outside Enterprise palette", f"Detected font color: {color}", "Use colors from the Enterprise color palette.", 0.55, shape))
    return issues[:6]


def check_footnotes(slide, slide_num: int) -> list[Issue]:
    text = slide_text(slide)
    # Find references like (1), (2), excluding dates and numbers in tables is imperfect.
    markers = set(re.findall(r"\((\d+)\)", text))
    issues = []
    if markers:
        # crude check: footnote entries often contain the marker plus source/footer text near bottom; if only one instance of marker, flag review.
        for m in sorted(markers):
            if len(re.findall(rf"\({m}\)", text)) < 2 and not re.search(rf"\b{m}\s+Source|Source.*\({m}\)", text, flags=re.I):
                issues.append(make_issue(slide_num, "FOOTNOTE_REVIEW", "Footnotes", "low", f"Footnote marker ({m}) may need review", f"Found marker ({m}) with no obvious matching note entry.", "Confirm every footnote marker has a matching footnote entry on the same slide.", 0.45))
    return issues[:4]


def check_notes_expected(slide, slide_num: int) -> list[Issue]:
    notes = extract_notes(slide)
    issues = []
    # Do not treat notes as app findings for production decks; surface them as evaluation notes only.
    for note in notes:
        if any(k in note.lower() for k in ["confidentiality", "privacy statement", "table", "title", "key", "legend", "bullet", "footnote"]):
            issues.append(make_issue(slide_num, "ANNOTATED_NOTE", "Annotated deck note", "low", "Annotated deck note extracted", note, "Use this note to compare whether the rule engine is catching expected issues. Not shown as a production finding by default.", 0.99))
    return issues


def analyze_presentation(path: str | Path, include_annotated_notes: bool = True) -> AuditReport:
    pptx_path = Path(path)
    prs = load_presentation(pptx_path)
    issues: list[Issue] = []
    summaries: list[SlideSummary] = []
    notes_count = 0

    for idx, slide in enumerate(prs.slides, start=1):
        slide_issues = []
        slide_issues.extend(check_title(slide, idx))
        slide_issues.extend(check_confidentiality(slide, idx))
        slide_issues.extend(check_title_slide(slide, idx))
        slide_issues.extend(check_layout_bounds(slide, idx))
        slide_issues.extend(check_bullets(slide, idx))
        slide_issues.extend(check_terminology(slide, idx))
        slide_issues.extend(check_tables(slide, idx))
        slide_issues.extend(check_colors(slide, idx))
        slide_issues.extend(check_footnotes(slide, idx))

        notes = extract_notes(slide)
        notes_count += len(notes)
        if include_annotated_notes and notes:
            # Store notes as low-severity reference issues only when user chooses annotated-note mode.
            slide_issues.extend(check_notes_expected(slide, idx))

        # Noise control: keep high/medium first; low confidence later.
        slide_issues.sort(key=lambda x: ({"high": 0, "medium": 1, "low": 2}[x.severity], -x.confidence, x.rule_id))
        issues.extend(slide_issues)

        title_shape = get_title_shape(slide)
        title = normalize_text(shape_text(title_shape)) if title_shape is not None else f"Slide {idx}"
        summaries.append(SlideSummary(
            slide_number=idx,
            title=title[:160] or f"Slide {idx}",
            issue_count=len(slide_issues),
            high_count=sum(1 for i in slide_issues if i.severity == "high"),
            medium_count=sum(1 for i in slide_issues if i.severity == "medium"),
            low_count=sum(1 for i in slide_issues if i.severity == "low"),
            notes=notes,
        ))

    return AuditReport(
        file_name=pptx_path.name,
        slide_count=len(prs.slides),
        issues=issues,
        slide_summaries=summaries,
        annotated_notes_found=notes_count,
    )
